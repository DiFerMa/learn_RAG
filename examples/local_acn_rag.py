import os
import json
import pprint
from pptx import Presentation
from collections import defaultdict
from datetime import datetime
import unicodedata
import re

input_folder = r"C:\Users\d.fernandez.macias\Desktop\temp"
zones_path = "zones.json"

def save_results(parsed_results):
    with open("rag_dataset.json", "w", encoding="utf-8") as f:
        json.dump(parsed_results, f, ensure_ascii=False, indent=2)

def join_content(content):
    relevant_keys = ["projects","profile","industry_experience","expertise"]
    joined_strings = ""
    for entry in relevant_keys:
        value = content.get(entry, [''])
        if value:
            joined_strings += f"{entry}: {value[0].strip()}\n"
    return joined_strings.strip()

def clean_text(text):
    if not text:
        return ""
    text = unicodedata.normalize("NFKC", text)
    text = text.replace('\xa0', ' ')
    # Remove control characters and non-printable characters (except newline and tab)
    text = re.sub(r'[^\x20-\x7E\n\t°€¥£±µ]', '', text)
    text = re.sub(r'\s+', ' ', text).strip()
    return text

def load_zones(zones_path):
    with open(zones_path, 'r') as f:
        return json.load(f)

def classify_shape_by_position(shape, zones):
    for section, bounds in zones.items():
        top_min, top_max = bounds.get("top", [None, None])
        left_min, left_max = bounds.get("left", [None, None])
        if top_min is not None and not (top_min <= shape.top <= top_max):
            continue
        if left_min is not None and not (left_min <= shape.left <= left_max):
            continue
        return section
    return "unknown"

def parse_slide(pptx_path, zones):
    prs = Presentation(pptx_path)
    slide = prs.slides[0]
    content = defaultdict(list)
    for shape in slide.shapes:
        if shape.has_text_frame:
            section = classify_shape_by_position(shape, zones)
            cleaned = clean_text(shape.text)
            content[section].append(cleaned)
    return dict(content)

def get_file_metadata(filepath):
    stat = os.stat(filepath)
    return {
        "file_name": os.path.basename(filepath),
        "last_modified": datetime.fromtimestamp(stat.st_mtime).isoformat()
    }

def process_pptx_folder(folder_path, zones_path):
    zones = load_zones(zones_path)
    results = []
    for file in os.listdir(folder_path):
        if not file.lower().endswith(".pptx"):
            continue
        if file.startswith("~$"):  # skip temp files
            continue
        pptx_file_path = os.path.join(folder_path, file)
        metadata = get_file_metadata(pptx_file_path)
        try:
            content = parse_slide(pptx_file_path, zones)
            prepared_content = join_content(content)
            metadata["email"] = content.get("email", [""])[0]
            metadata["phone"] = content.get("phone", [""])[0]
            metadata["location"] = content.get("location", [""])[0]
            metadata["title"] = content.get("title", [""])[0]
            results.append({
                "metadata": metadata,
                "content": prepared_content
            })
        except Exception as e:
            print(f"Error parsing {file}: {e}")

    return results
if __name__ == "__main__":
    parsed_results = process_pptx_folder(input_folder, zones_path)
    save_results(parsed_results)
